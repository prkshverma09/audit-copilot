import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 Widescreen (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Palette
    C_BG = RGBColor(10, 15, 29)       # #0a0f1d
    C_PANEL = RGBColor(17, 24, 43)    # #11182b
    C_CARD = RGBColor(22, 32, 57)     # #162039
    C_BORDER = RGBColor(38, 52, 85)   # #263455
    C_WHITE = RGBColor(255, 255, 255)
    C_SLATE = RGBColor(148, 163, 184) # #94a3b8
    C_SKY = RGBColor(56, 189, 248)    # #38bdf8
    C_EMERALD = RGBColor(16, 185, 129)# #10b981
    C_AMBER = RGBColor(245, 158, 11)  # #f59e0b
    C_ROSE = RGBColor(244, 63, 94)    # #f43f5e
    
    ASSETS_DIR = os.path.join(os.path.dirname(__file__), "assets")
    concept_img = os.path.join(ASSETS_DIR, "concept.jpg")
    arch_img = os.path.join(ASSETS_DIR, "architecture.jpg")
    app_img = os.path.join(ASSETS_DIR, "app_screenshot.png")
    
    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_BG
        bg.line.color.rgb = C_BG
        return bg

    def add_header(slide, tag_text, title_text, subtitle_text):
        # Tag Badge
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.35))
        tf = tag_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p_tag = tf.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.size = Pt(10)
        p_tag.font.bold = True
        p_tag.font.color.rgb = C_SKY
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.55))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = C_WHITE
        
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.35))
        tf_s = sub_box.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
        p_s = tf_s.paragraphs[0]
        p_s.text = subtitle_text
        p_s.font.size = Pt(13)
        p_s.font.color.rgb = C_SLATE

    # ==========================================
    # SLIDE 1: PROBLEM & ORIGIN
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)
    add_header(s1, "Ylookup x Encode AI Hackathon | Product Track", 
               "The Problem: Manual Reconciliation in Private Capital",
               "100+ quarterly auditor hours lost cross-referencing spreadsheets against PDF statements")
    
    # Left: 3 Concise Problem Cards
    cards_data = [
        ("01 • NO DOCUMENT PROVENANCE", "Formulas like C4 + D5 = C6 reference ledger cells with zero verifiable links to source bank PDFs.", C_ROSE),
        ("02 • HIGH-RISK MANUAL ENTRY", "Auditors manually re-type balances across multiple funds. A single typo triggers days of reconciliation delays.", C_AMBER),
        ("03 • HIDDEN EXCEPTIONS", "Unallocated deposits and timing discrepancies stay buried in ledger rows until year-end fire drills.", C_SKY)
    ]
    
    y_start = 1.9
    for i, (head, desc, color) in enumerate(cards_data):
        y = y_start + (i * 1.65)
        card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(6.0), Inches(1.45))
        card.fill.solid()
        card.fill.fore_color.rgb = C_PANEL
        card.line.color.rgb = C_BORDER
        card.line.width = Pt(1)
        
        tb = s1.shapes.add_textbox(Inches(1.05), Inches(y + 0.18), Inches(5.5), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p0 = tf.paragraphs[0]
        p0.text = head
        p0.font.size = Pt(12)
        p0.font.bold = True
        p0.font.color.rgb = color
        p0.space_after = Pt(4)
        
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(11)
        p1.font.color.rgb = C_SLATE

    # Right: Quote & Stat Pills
    q_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(1.9), Inches(5.4), Inches(4.75))
    q_card.fill.solid()
    q_card.fill.fore_color.rgb = C_PANEL
    q_card.line.color.rgb = C_BORDER
    q_card.line.width = Pt(1)
    
    q_tb = s1.shapes.add_textbox(Inches(7.45), Inches(2.2), Inches(4.7), Inches(4.2))
    q_tf = q_tb.text_frame
    q_tf.word_wrap = True
    q_tf.margin_left = q_tf.margin_top = q_tf.margin_right = q_tf.margin_bottom = 0
    
    qp0 = q_tf.paragraphs[0]
    qp0.text = "ORIGIN: FUND MANAGER INTERVIEWS"
    qp0.font.size = Pt(11)
    qp0.font.bold = True
    qp0.font.color.rgb = C_SKY
    qp0.space_after = Pt(14)
    
    qp1 = q_tf.add_paragraph()
    qp1.text = '“Auditors spend up to 70% of quarterly reporting manually matching Excel numbers against PDF bank statements. It is slow, disconnected, and error-prone.”'
    qp1.font.size = Pt(15)
    qp1.font.italic = True
    qp1.font.color.rgb = C_WHITE
    qp1.space_after = Pt(24)
    
    # 3 Stats in PPT
    stats = [
        ("70%", "MANUAL TIME"),
        ("2-4 WKS", "NAV DELAY"),
        ("100%", "PROOF REQUIRED")
    ]
    for val, lbl in stats:
        p_stat = q_tf.add_paragraph()
        p_stat.text = f"• {val} — {lbl}"
        p_stat.font.size = Pt(12)
        p_stat.font.bold = True
        p_stat.font.color.rgb = C_EMERALD
        p_stat.space_after = Pt(6)

    # ==========================================
    # SLIDE 2: SOLUTION (X-RAY AUDIT COPILOT)
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "Gemini 2.5 Multimodal AI | Product Solution",
               "The Solution: X-Ray Audit Copilot",
               "Gemini-powered financial lineage & sub-millimeter character-grounded statement matching")
    
    p_data = [
        ("Engine 1 • Gemini 2.5 Multimodal Lineage", "Google Gemini multimodal understanding extracts statement metadata, accounts, and verbatim citation quotes into structured JSON.", C_SKY),
        ("Engine 2 • Automated Math Verification", "Python AST engine recalculates vertical ledger sums (C4 + D5 = C6) and cross-fund balances with zero delta.", C_EMERALD),
        ("Engine 3 • Grounded Statement Matching", "Extracts exact bounding boxes (bbox) from PDF statements for 100% verified, non-hallucinatory citations.", C_AMBER)
    ]
    
    for i, (head, desc, color) in enumerate(p_data):
        y = 1.9 + (i * 1.65)
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(5.5), Inches(1.45))
        card.fill.solid()
        card.fill.fore_color.rgb = C_PANEL
        card.line.color.rgb = C_BORDER
        card.line.width = Pt(1)
        
        tb = s2.shapes.add_textbox(Inches(1.05), Inches(y + 0.18), Inches(5.0), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p0 = tf.paragraphs[0]
        p0.text = head
        p0.font.size = Pt(12)
        p0.font.bold = True
        p0.font.color.rgb = color
        p0.space_after = Pt(4)
        
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(11)
        p1.font.color.rgb = C_SLATE

    if os.path.exists(concept_img):
        s2.shapes.add_picture(concept_img, Inches(6.6), Inches(1.9), width=Inches(5.9))

    # ==========================================
    # SLIDE 3: HIGH LEVEL ARCHITECTURE
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "System Design | Zero-Hallucination Pipeline",
               "High-Level System Architecture",
               "Deterministic 4-stage pipeline connecting multi-fund sheets to source bank statements")
    
    if os.path.exists(arch_img):
        s3.shapes.add_picture(arch_img, Inches(0.8), Inches(1.85), width=Inches(7.8))
    
    arch_card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.85), Inches(3.6), Inches(4.85))
    arch_card.fill.solid()
    arch_card.fill.fore_color.rgb = C_PANEL
    arch_card.line.color.rgb = C_BORDER
    arch_card.line.width = Pt(1)
    
    a_tb = s3.shapes.add_textbox(Inches(9.15), Inches(2.05), Inches(3.1), Inches(4.4))
    a_tf = a_tb.text_frame
    a_tf.word_wrap = True
    a_tf.margin_left = a_tf.margin_top = a_tf.margin_right = a_tf.margin_bottom = 0
    
    ap0 = a_tf.paragraphs[0]
    ap0.text = "CORE GUARANTEES"
    ap0.font.size = Pt(11)
    ap0.font.bold = True
    ap0.font.color.rgb = C_SKY
    ap0.space_after = Pt(12)
    
    highlights = [
        ("Gemini 2.5 Flash & Pro", "Flash for rapid PDF classification; Pro for multimodal financial extraction."),
        ("Sub-Millimeter BBox", "Direct character coordinate extraction via PyMuPDF and pdfplumber."),
        ("Deterministic AST Math", "Python AST formula evaluation; zero LLM math guessing or hallucination."),
        ("Suspense Alerting", "Detects unallocated deposits (e.g. €45,200) with review badges.")
    ]
    
    for title, detail in highlights:
        pt = a_tf.add_paragraph()
        pt.text = "• " + title
        pt.font.size = Pt(11)
        pt.font.bold = True
        pt.font.color.rgb = C_WHITE
        
        pd = a_tf.add_paragraph()
        pd.text = "   " + detail
        pd.font.size = Pt(10)
        pd.font.color.rgb = C_SLATE
        pd.space_after = Pt(8)

    # ==========================================
    # SLIDE 4: TECHNOLOGIES USED & DEPLOYMENT
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "Engineering & Delivery | Production Ready",
               "Technology Stack & Turnkey Delivery",
               "Modern, decoupled architecture with single-command Docker deployment")
    
    tech_cards = [
        ("GOOGLE GEMINI 2.5 (FLASH & PRO)", "Official google-genai SDK. Flash for fast PDF classification; Pro for structured financial entity & quote extraction.", C_SKY),
        ("FASTAPI & PYTHON 3.11", "Asynchronous REST API, formula AST dependency solver, PyMuPDF coordinate geometry, sub-50ms latency.", C_EMERALD),
        ("NEXT.JS 14 & FORTUNESHEET", "React 18 App Router, Tailwind CSS, virtualized spreadsheet canvas, and synchronized PDF.js split-viewer.", C_SKY),
        ("TURNKEY DOCKER", "Single-command containerization (./run_docker.sh) running full stack with optional GEMINI_API_KEY or offline heuristic fallback.", C_AMBER)
    ]
    
    grid_coords = [
        (Inches(0.8), Inches(1.9)),
        (Inches(6.8), Inches(1.9)),
        (Inches(0.8), Inches(4.35)),
        (Inches(6.8), Inches(4.35)),
    ]
    
    for (x, y), (head, desc, color) in zip(grid_coords, tech_cards):
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.7), Inches(2.15))
        card.fill.solid()
        card.fill.fore_color.rgb = C_PANEL
        card.line.color.rgb = C_BORDER
        card.line.width = Pt(1)
        
        tb = s4.shapes.add_textbox(x + Inches(0.25), y + Inches(0.25), Inches(5.2), Inches(1.65))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p0 = tf.paragraphs[0]
        p0.text = head
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = color
        p0.space_after = Pt(6)
        
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(11)
        p1.font.color.rgb = C_SLATE
        p1.space_after = Pt(8)
    
    output_path = os.path.join(os.path.dirname(__file__), "X-Ray_Audit_Copilot_Demo_Slides.pptx")
    prs.save(output_path)
    print(f"Presentation successfully saved to: {output_path}")

if __name__ == "__main__":
    create_deck()
